import os
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from io import BytesIO

class DocumentProcessor:
    @staticmethod
    def _sanitize_text(text: str) -> str:
        """Remove NUL bytes and other problematic characters from text"""
        if not text:
            return ""
        # Remove NUL bytes which cause PostgreSQL errors
        return text.replace('\x00', '')

    @staticmethod
    def extract_text(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            text = DocumentProcessor._extract_from_pdf(file_path)
        elif ext == '.docx':
            text = DocumentProcessor._extract_from_docx(file_path)
        elif ext == '.pptx':
            text = DocumentProcessor._extract_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
        return DocumentProcessor._sanitize_text(text)

    @staticmethod
    def extract_text_from_bytes(file_bytes: bytes, filename: str):
        ext = os.path.splitext(filename)[1].lower()
        bio = BytesIO(file_bytes)
        if ext == '.pdf':
            text = DocumentProcessor._extract_pdf_bytes(bio)
        elif ext == '.docx':
            text = DocumentProcessor._extract_docx_bytes(bio)
        elif ext == '.pptx':
            text = DocumentProcessor._extract_pptx_bytes(bio)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
        return DocumentProcessor._sanitize_text(text)

    @staticmethod
    def _extract_from_pdf(file_path):
        from app.services.ai_service import AIService
        text = ""
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
                # Multimodal Image Captioning
                try:
                    if hasattr(page, 'images') and page.images:
                        for img in page.images:
                             # img.data contains the bytes
                             caption = AIService.generate_image_caption(img.data)
                             text += caption + "\n"
                except Exception as e:
                     # Non-blocking error for image extraction
                     print(f"Warning: Failed to extract images from PDF page: {e}")
        return text

    @staticmethod
    def _extract_pdf_bytes(bio: BytesIO):
        from app.services.ai_service import AIService
        text = ""
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text += page.extract_text() + "\n"
            # Multimodal Image Captioning
            try:
                if hasattr(page, 'images') and page.images:
                    for img in page.images:
                         caption = AIService.generate_image_caption(img.data)
                         text += caption + "\n"
            except Exception as e:
                 print(f"Warning: Failed to extract images from PDF page: {e}")
        return text

    @staticmethod
    def _extract_from_docx(file_path):
        doc = DocxDocument(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _extract_docx_bytes(bio: BytesIO):
        doc = DocxDocument(bio)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _extract_from_pptx(file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    def _extract_pptx_bytes(bio: BytesIO):
        prs = Presentation(bio)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    def chunk_text(text, chunk_size=512, overlap=50):
        # Simple character/word based chunking for MVP
        # Ideally use tiktoken or similar for token-based
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            chunks.append(chunk)
            
        return chunks
